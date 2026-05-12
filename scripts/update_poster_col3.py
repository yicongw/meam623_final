#!/usr/bin/env python3
"""Update Col3 (Swapping DS Planners) — unified style.

Same shared spec as Col1/Col2 (scripts/poster_style.py).
"""

from __future__ import annotations
import os, shutil, sys

from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE

_HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, _HERE)
import poster_style as S
from poster_style import NAVY, GREEN, DARK, BODY, MID
from pptx.util import Pt

POSTER = r"C:\Users\wangy\Desktop\meam623_poster.pptx"
BACKUP = r"C:\Users\wangy\Desktop\meam623_poster_backup_before_col3.pptx"
FIG    = r"C:\meam623_finalproj\output\fig_poster_col3_lpvds_only.png"


def _clear(tf):
    p = tf.paragraphs[0]._p
    for sib in list(p.itersiblings()):
        sib.getparent().remove(sib)
    for r in list(p.findall(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}r")):
        p.remove(r)


def _add(tf, runs, *, align=None, first=False):
    para = tf.paragraphs[0] if first else tf.add_paragraph()
    if align is not None:
        para.alignment = align
    for spec in runs:
        run = para.add_run()
        run.text = spec["text"]
        f = run.font
        f.name = spec.get("name", S.FONT)
        f.size = Pt(spec["size_pt"])
        f.bold = spec.get("bold", False)
        f.italic = spec.get("italic", False)
        if spec.get("color") is not None:
            f.color.rgb = spec["color"]
        if spec.get("baseline") is not None:
            run._r.get_or_add_rPr().set("baseline", str(spec["baseline"]))


# Inline-math fragment helpers (LaTeX-style).
SUB = -25000
SUP =  30000


def _math(text, size_pt, *, color, bold=False):
    return {"text": text, "size_pt": size_pt, "italic": True,
            "bold": bold, "color": color}


def _sub(text, size_pt, *, color, bold=False):
    return {"text": text, "size_pt": size_pt, "italic": True,
            "bold": bold, "color": color, "baseline": SUB}


def _sup(text, size_pt, *, color, bold=False):
    return {"text": text, "size_pt": size_pt, "italic": False,
            "bold": bold, "color": color, "baseline": SUP}


def _plain(text, size_pt, *, color, bold=False):
    return {"text": text, "size_pt": size_pt, "italic": False,
            "bold": bold, "color": color}


def write_col3_title(tf):
    _clear(tf)
    _add(tf, [
        {"text": "3. Swapping DS Planners", "size_pt": S.TITLE,
         "bold": True, "color": NAVY},
        {"text": "  \u2713 Done", "size_pt": S.DONE,
         "bold": False, "color": GREEN},
    ], align=PP_ALIGN.LEFT, first=True)
    _add(tf, [
        {"text": "Drop in any GAS dynamical system as the high-level "
                 "planner f(x); the QP and safety layer are untouched. "
                 "Tested on a 2-obstacle reach task in PyBullet.",
         "size_pt": S.BODY_PT, "bold": False, "color": BODY},
    ])


def write_col3_box(tf):
    """LaTeX-style equation: f(x) = Σₖ γₖ(x)·(Aₖ x + bₖ),  Aₖ + Aₖᵀ ≺ 0."""
    _clear(tf)
    em = S.EQ_MAIN

    _add(tf, [
        _math ("f", em, color=NAVY, bold=True),
        _plain("(", em, color=NAVY, bold=True),
        _math ("x", em, color=NAVY, bold=True),
        _plain(") = \u03a3", em, color=NAVY, bold=True),
        _sub  ("k", em, color=NAVY, bold=True),
        _plain(" ", em, color=NAVY, bold=True),
        _math ("\u03b3", em, color=NAVY, bold=True),          # γ
        _sub  ("k", em, color=NAVY, bold=True),
        _plain("(", em, color=NAVY, bold=True),
        _math ("x", em, color=NAVY, bold=True),
        _plain(") \u00b7 (", em, color=NAVY, bold=True),
        _math ("A", em, color=NAVY, bold=True),
        _sub  ("k", em, color=NAVY, bold=True),
        _plain(" ", em, color=NAVY, bold=True),
        _math ("x", em, color=NAVY, bold=True),
        _plain(" + ", em, color=NAVY, bold=True),
        _math ("b", em, color=NAVY, bold=True),
        _sub  ("k", em, color=NAVY, bold=True),
        _plain(",   ", em, color=NAVY, bold=True),
        _math ("A", em, color=NAVY, bold=True),
        _sub  ("k", em, color=NAVY, bold=True),
        _plain(" + ", em, color=NAVY, bold=True),
        _math ("A", em, color=NAVY, bold=True),
        _sub  ("k", em, color=NAVY, bold=True),
        _sup  ("\u22a4", em, color=NAVY, bold=True),          # ⊤
        _plain(" \u227a 0", em, color=NAVY, bold=True),
    ], align=PP_ALIGN.CENTER, first=True)

    es = S.EQ_SUB
    _add(tf, [
        _plain("Hand-drawn LPV-DS  \u00b7  ", es, color=MID),
        _math ("K", es, color=MID),
        _plain(" = 6 components  \u00b7  Hurwitz-projected", es, color=MID),
    ], align=PP_ALIGN.CENTER)


def write_col3_plan(tf):
    _clear(tf)
    bs = S.BULL_BODY
    ph = S.PLAN_HEAD

    # Headline: "Same QP & safety layer — only f(x) is swapped."
    _add(tf, [
        _plain("Same QP & safety layer  \u2014  only ", ph,
               color=NAVY, bold=True),
        _math ("f", ph, color=NAVY, bold=True),
        _plain("(", ph, color=NAVY, bold=True),
        _math ("x", ph, color=NAVY, bold=True),
        _plain(") is swapped.", ph, color=NAVY, bold=True),
    ], align=PP_ALIGN.LEFT, first=True)

    bullets = [
        # bullet 1: min d_obs = 7.9 cm,  min Γ = 15.9
        ("\u2713  Safety preserved", [
            _plain("    min ", bs, color=BODY),
            _math ("d", bs, color=BODY),
            _sub  ("obs", bs, color=BODY),
            _plain(" = 7.9 cm,  min ", bs, color=BODY),
            _math ("\u0393", bs, color=BODY),                 # Γ
            _plain(" = 15.9  (no obstacle contact, "
                   "no self-collision).", bs, color=BODY),
        ]),
        # bullet 2: S(t) = T_kin + λ V_pot
        ("\u2713  Empirical passivity", [
            _plain("    ", bs, color=BODY),
            _math ("S", bs, color=BODY),
            _plain("(", bs, color=BODY),
            _math ("t", bs, color=BODY),
            _plain(") = ", bs, color=BODY),
            _math ("T", bs, color=BODY),
            _sub  ("kin", bs, color=BODY),
            _plain(" + ", bs, color=BODY),
            _math ("\u03bb", bs, color=BODY),                 # λ
            _plain(" ", bs, color=BODY),
            _math ("V", bs, color=BODY),
            _sub  ("pot", bs, color=BODY),
            _plain(" stays bounded throughout the run.",
                   bs, color=BODY),
        ]),
        # bullet 3: no math
        ("\u2713  Hand-drawn LPV-DS converges", [
            _plain("    Final EE-to-target error 14.7 cm under "
                   "viability constraints.", bs, color=BODY),
        ]),
    ]
    for header, body_runs in bullets:
        _add(tf, [{"text": header, "size_pt": S.BULL_HEAD,
                   "bold": True, "color": DARK}])
        _add(tf, body_runs)

    _add(tf, [
        {"text": "Planner-agnostic: any GAS DS plugs in.",
         "size_pt": S.TAKEAWAY, "bold": True, "color": NAVY},
    ])


def replace_picture(slide, target_name, new_image_path,
                    *, expected_left_in=None, expected_top_in=None, tol_in=0.5):
    target = None
    for shape in slide.shapes:
        if shape.name == target_name:
            target = shape; break
    if target is None and expected_left_in is not None:
        EMU = 914400
        best, best_d = None, float("inf")
        for shape in slide.shapes:
            if shape.shape_type != 13:
                continue
            d = abs(shape.left/EMU - expected_left_in) + \
                abs(shape.top /EMU - expected_top_in)
            if d < best_d:
                best, best_d = shape, d
        if best is not None and best_d < tol_in * 2:
            print(f"  fallback: matched {target_name!r} by position "
                  f"(was {best.name!r}, dist={best_d:.2f}in)")
            target = best
    if target is None:
        raise RuntimeError(f"Picture {target_name!r} not found")
    left, top, w, h = target.left, target.top, target.width, target.height
    sp = target._element
    sp.getparent().remove(sp)
    pic = slide.shapes.add_picture(new_image_path, left, top, width=w, height=h)
    pic.name = target_name
    return pic


def main():
    if not os.path.exists(BACKUP):
        shutil.copy2(POSTER, BACKUP)
        print(f"backup -> {BACKUP}")
    p = Presentation(POSTER)
    slide = p.slides[0]
    by_name = {sh.name: sh for sh in slide.shapes}
    for need in ("Col3Title", "Col3Box", "Col3Plan"):
        if need not in by_name:
            raise RuntimeError(f"missing shape: {need}")
    for nm in ("Col3Title", "Col3Box", "Col3Plan"):
        tf = by_name[nm].text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.word_wrap = True
    write_col3_title(by_name["Col3Title"].text_frame)
    write_col3_box  (by_name["Col3Box"  ].text_frame)
    write_col3_plan (by_name["Col3Plan" ].text_frame)
    replace_picture (slide, "Col3Fig", FIG,
                     expected_left_in=22.20, expected_top_in=34.65)
    p.save(POSTER)
    print(f"saved -> {POSTER}")


if __name__ == "__main__":
    main()
