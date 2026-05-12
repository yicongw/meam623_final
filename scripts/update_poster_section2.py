"""Update Section 2 of meam623_poster.pptx with the dynamics-mismatch results.

Mirrors the layout used for Section 3:
  - Col2Title : 'Done' badge (green) + corrected one-line summary.
  - Col2Equation: compact box stating the perturbed QP and the swept range.
  - Inserts poster_fig2.png between Col2Equation and Col2Plan.
  - Col2Plan : actual results / take-aways (replaces 'To be shown').
"""

from __future__ import annotations
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

POSTER = Path(r"C:/Users/wangy/Desktop/meam623_poster.pptx")
FIG    = Path(r"C:/Users/wangy/Desktop/poster_fig2.png")
BACKUP = POSTER.with_suffix(".section2_backup.pptx")

NAVY   = RGBColor(0x01, 0x1F, 0x5B)
GREEN  = RGBColor(0x2E, 0x7D, 0x32)
BODY   = RGBColor(0x33, 0x33, 0x33)
MUTED  = RGBColor(0x55, 0x55, 0x55)


def set_run(run, text, *, font="Arial", size=None, bold=None, color=None):
    run.text = text
    run.font.name = font
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def clear_text_frame(tf):
    p_elems = tf._txBody.findall(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}p"
    )
    for p in p_elems[1:]:
        tf._txBody.remove(p)
    first = tf.paragraphs[0]
    for r in list(first.runs):
        first._p.remove(r._r)


def add_para(tf, *, alignment=None):
    p = tf.add_paragraph()
    if alignment is not None:
        p.alignment = alignment
    return p


def first_para(tf, *, alignment=None):
    p = tf.paragraphs[0]
    if alignment is not None:
        p.alignment = alignment
    return p


def main():
    if not POSTER.exists():
        raise FileNotFoundError(POSTER)
    if not FIG.exists():
        raise FileNotFoundError(FIG)

    if not BACKUP.exists():
        shutil.copy2(POSTER, BACKUP)
        print(f"[backup] -> {BACKUP}")

    prs   = Presentation(str(POSTER))
    slide = prs.slides[0]
    by    = {sh.name: sh for sh in slide.shapes}
    title = by["Col2Title"]
    eq    = by["Col2Equation"]
    plan  = by["Col2Plan"]

    # -------------------------------------------------------------------- #
    # 1. Col2Title: 'Done' badge + one-line summary of the experiment.     #
    # -------------------------------------------------------------------- #
    tf = title.text_frame
    clear_text_frame(tf)

    p0 = first_para(tf)
    r = p0.add_run()
    set_run(r, "2. Model-Mismatch Robustness", size=40, bold=True, color=NAVY)
    r = p0.add_run()
    set_run(r, "  \u2713 Done", size=28, bold=False, color=GREEN)

    p1 = add_para(tf)
    r = p1.add_run()
    set_run(
        r,
        "Aggressive scenario (target near base, gain 300, no obstacles, no "
        "reactive evasion). QP sees \u03B1\u00B7M\u209C\u1D63\u1D64\u2091; "
        "swept \u03B1 \u2208 [0.5, 1.5] to compare with a CFC ablation.",
        size=32, bold=False, color=BODY,
    )

    # -------------------------------------------------------------------- #
    # 2. Col2Equation: compact box.                                         #
    # -------------------------------------------------------------------- #
    eq.top    = Inches(33.20)
    eq.height = Inches(1.30)

    tf = eq.text_frame
    clear_text_frame(tf)

    p0 = first_para(tf, alignment=PP_ALIGN.CENTER)
    r = p0.add_run()
    set_run(
        r,
        "VPP-TC:  \u00C2(\u03B1\u00B7M)\u00B7\u03C4 \u2265 b\u0302 \u2212 \u03B4  +  "
        "viability  +  \u0393   vs.   CFC: drop both",
        size=30, bold=True, color=NAVY,
    )

    p1 = add_para(tf, alignment=PP_ALIGN.CENTER)
    r = p1.add_run()
    set_run(
        r,
        "Same QP backbone, same dynamics-mismatch sweep \u03B1 \u2208 [0.5, 1.5]",
        size=22, bold=False, color=MUTED,
    )

    # -------------------------------------------------------------------- #
    # 3. Insert poster_fig2.png between Col2Equation and Col2Plan.          #
    # -------------------------------------------------------------------- #
    img_left   = title.left
    img_width  = title.width
    img_height = Emu(int(img_width * 6 / 15))   # 15:6 aspect
    img_top    = Inches(34.65)

    for sh in list(slide.shapes):
        if sh.name == "Col2Fig":
            sp = sh._element
            sp.getparent().remove(sp)

    pic = slide.shapes.add_picture(
        str(FIG),
        left=img_left, top=img_top,
        width=img_width, height=img_height,
    )
    pic.name = "Col2Fig"

    # -------------------------------------------------------------------- #
    # 4. Col2Plan: actual results / take-aways.                             #
    # -------------------------------------------------------------------- #
    plan_top    = img_top + img_height + Inches(0.20)   # ~38.70 in
    plan_height = Inches(43.30) - plan_top              # stop above footer
    plan.top    = plan_top
    plan.height = plan_height

    tf = plan.text_frame
    clear_text_frame(tf)

    p0 = first_para(tf)
    r = p0.add_run()
    set_run(
        r,
        "VPP-TC:  0 / 7 collisions     CFC:  2 / 7 collisions",
        size=28, bold=True, color=NAVY,
    )

    bullets = [
        ("\u2713  VPP-TC stays safe across the full sweep",
         "min d\u209B\u1D9C \u2208 [2.7, 3.8] cm; \u0393 dips but never goes negative"),
        ("\u2717  CFC self-collides at \u03B1 = 1.00 and 1.15",
         "even with the perfect inertia model: viability + \u0393 are essential"),
        ("\u2713  VPP-TC reaches the target 7 / 7;  CFC only 5 / 7",
         "(worst-case final \u2016x \u2212 x*\u2016 = 89 mm for CFC, < 0.1 mm for VPP-TC)"),
    ]
    for head, sub in bullets:
        p = add_para(tf)
        r = p.add_run()
        set_run(r, head, size=26, bold=True, color=BODY)
        p2 = add_para(tf)
        r = p2.add_run()
        set_run(r, "    " + sub, size=20, bold=False, color=MUTED)

    p_take = add_para(tf)
    r = p_take.add_run()
    set_run(
        r,
        "Viability + \u0393 are what buy robustness -- the QP backbone alone is not enough.",
        size=26, bold=True, color=NAVY,
    )

    prs.save(str(POSTER))
    print(f"[save] -> {POSTER}")


if __name__ == "__main__":
    main()
