#!/usr/bin/env python3
"""Restyle Col1 (Multi-Platform Porting) to match Col2/Col3 unified style
AND fix the layout — close the dead-space gap between the two Franka figures
and the caption, bottom-align with Col2/Col3, and extend the separator
lines to the same baseline.

In-place edit of C:\\Users\\wangy\\Desktop\\meam623_poster.pptx.
"""

from __future__ import annotations
import os, shutil, sys

from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_ANCHOR
from pptx.util import Pt, Emu

_HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, _HERE)
import poster_style as S
from poster_style import NAVY, GREEN, DARK, BODY, MID

POSTER = r"C:\Users\wangy\Desktop\meam623_poster.pptx"
BACKUP = r"C:\Users\wangy\Desktop\meam623_poster_backup_before_col1.pptx"

EMU = 914400

# Target geometry (inches), parallel 4-row structure with Col2/Col3:
#   title   30.62 -> 33.03    (unchanged)
#   setup   33.20 -> 34.50    (NEW shape Col1Setup, parallel to Col2Equation)
#   figs    34.65 -> 38.79    (preserves original aspect, H=4.14)
#   caption 38.92 -> 43.70    (matches Col2Plan top + height)
TARGET_SETUP_TOP_IN   = 33.40
TARGET_SETUP_H_IN     = 1.70    # bigger to fit EQ_MAIN 28 + EQ_SUB 22
TARGET_FIG_TOP_IN     = 35.20   # pushed down to make room for taller setup row
TARGET_FIG_HEIGHT_IN  = 4.14    # Col1: preserve native aspect ratio
TARGET_CAP_TOP_IN     = 40.15   # plan-row top, accommodates taller Col3Fig below
TARGET_CAP_HEIGHT_IN  = 5.95    # fill all the way to the results-card bottom
SEP_NEW_BOTTOM_IN     = 46.10   # results-card bottom is 46.20; leave 0.10 pad
COL1_LEFT_IN          = 1.20
COL1_WIDTH_IN         = 9.95
# Col3Fig regenerated at figsize=(10, 5.0) -> aspect 2.0; fill the full
# Col3 width and accept a small gap below the (smaller) Col1/Col2 figures.
COL3_FIG_LEFT_IN      = 22.20
COL3_FIG_WIDTH_IN     = 9.62
COL3_FIG_HEIGHT_IN    = 4.81    # 9.62 / 2.0 aspect from the matplotlib fig
# Col2Fig: time-series plot at column width.  Force H so it never bleeds
# into the plan-row above (PowerPoint's round-trip sometimes auto-resizes
# pictures back to their native aspect, which would push H to ~6in).
COL2_FIG_LEFT_IN      = 11.63
COL2_FIG_WIDTH_IN     = 10.25
COL2_FIG_HEIGHT_IN    = 4.81    # match Col3Fig bottom for clean alignment

_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _strip_bullets(para):
    from lxml import etree
    pPr = para._p.find(f"{_A}pPr")
    if pPr is None:
        pPr = etree.SubElement(para._p, f"{_A}pPr")
        para._p.insert(0, pPr)
    for tag in ("buAutoNum", "buChar", "buNone", "buFont"):
        for el in pPr.findall(f"{_A}{tag}"):
            pPr.remove(el)
    etree.SubElement(pPr, f"{_A}buNone")


def _clear(tf):
    p = tf.paragraphs[0]._p
    for sib in list(p.itersiblings()):
        sib.getparent().remove(sib)
    for r in list(p.findall(f"{_A}r")):
        p.remove(r)


def _add(tf, runs, *, align=None, first=False, no_bullet=True):
    para = tf.paragraphs[0] if first else tf.add_paragraph()
    if align is not None:
        para.alignment = align
    if no_bullet:
        _strip_bullets(para)
    for spec in runs:
        run = para.add_run()
        run.text = spec["text"]
        f = run.font
        f.name = spec.get("name", S.FONT)
        f.size = Pt(spec["size_pt"])
        f.bold = spec.get("bold", False)
        f.italic = spec.get("italic", False)
        if spec.get("color") is not None:
            f.color.rgb = spec["color"]


def _ensure_col1_setup(slide, title):
    """Create or update a 'Col1Setup' textbox at T=33.20, H=1.30 to mirror
    the Col2Equation / Col3Box row.  Matches their fill (#F5F5F5) and
    border (NAVY 0.5pt).  Returns the shape."""
    from pptx.dml.color import RGBColor
    from pptx.util import Pt as _Pt
    EMU_F = 914400
    by_name = {sh.name: sh for sh in slide.shapes}
    setup = by_name.get("Col1Setup")
    if setup is None:
        setup = slide.shapes.add_textbox(
            title.left,
            Emu(int(TARGET_SETUP_TOP_IN * EMU_F)),
            title.width,
            Emu(int(TARGET_SETUP_H_IN * EMU_F)),
        )
        setup.name = "Col1Setup"
        print(f"  created Col1Setup at T={TARGET_SETUP_TOP_IN:.2f} "
              f"H={TARGET_SETUP_H_IN:.2f}")
    else:
        setup.left   = title.left
        setup.width  = title.width
        setup.top    = Emu(int(TARGET_SETUP_TOP_IN * EMU_F))
        setup.height = Emu(int(TARGET_SETUP_H_IN   * EMU_F))
    # Match Col2Equation / Col3Box styling
    setup.fill.solid()
    setup.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
    setup.line.color.rgb      = RGBColor(0x01, 0x1F, 0x5B)
    setup.line.width          = _Pt(0.5)
    # Critical: disable auto-shrink so the box stays at TARGET_SETUP_H_IN
    # instead of collapsing to fit the text (which was causing the gray fill
    # to only cover the top half).
    setup.text_frame.auto_size = MSO_AUTO_SIZE.NONE
    setup.text_frame.word_wrap = True
    # Vertically center the text inside the gray box so it visually matches
    # Col2Equation / Col3Box (which look centered).
    setup.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return setup


def write_col1_setup(tf):
    """One-line 'equation-row' equivalent for Col1: a setup tagline that
    visually balances Col2Equation and Col3Box."""
    _clear(tf)
    _add(tf, [
        {"text": "Same QP & safety layer  \u00b7  zero re-tuning",
         "size_pt": S.EQ_MAIN, "bold": True, "color": NAVY},
    ], align=PP_ALIGN.CENTER, first=True)
    _add(tf, [
        {"text": "Two robot platforms, two URDFs  \u2014  "
                 "one controller pipeline.",
         "size_pt": S.EQ_SUB, "bold": False, "color": MID},
    ], align=PP_ALIGN.CENTER)


def write_col1_title(tf):
    _clear(tf)
    _add(tf, [
        {"text": "1. Multi-Platform Porting", "size_pt": S.TITLE,
         "bold": True, "color": NAVY},
        {"text": "  \u2713 Done", "size_pt": S.DONE,
         "bold": False, "color": GREEN},
    ], align=PP_ALIGN.LEFT, first=True)
    _add(tf, [
        {"text": "Same VPP-TC pipeline, no controller retuning. "
                 "Tested on the Franka Panda dual-arm and the OpenArm "
                 "humanoid (both 14-DoF).",
         "size_pt": S.BODY_PT, "bold": False, "color": BODY},
    ])


def write_col1_caption(tf):
    """Plan-style caption — three takeaways padded to fill the larger box."""
    _clear(tf)

    _add(tf, [
        {"text": "Panda dual-arm  \u00b7  OpenArm humanoid  "
                 "(both 14-DoF)",
         "size_pt": S.PLAN_HEAD, "bold": True, "color": NAVY},
    ], align=PP_ALIGN.LEFT, first=True)

    bullets = [
        ("\u2713  Inter-arm collision avoidance",
         "    Both Panda arms track circular limit cycles while VPP-TC "
         "keeps them apart \u2014 no platform-specific tuning."),
        ("\u2713  OpenArm humanoid port",
         "    Same controller, new URDF and joint limits; viability "
         "and passivity hold without re-derivation."),
        ("\u2713  Zero re-tuning across platforms",
         "    Gains, weights, and safety margins identical to the "
         "single-arm Panda baseline."),
    ]
    for header, body in bullets:
        _add(tf, [{"text": header, "size_pt": S.BULL_HEAD,
                   "bold": True, "color": DARK}])
        _add(tf, [{"text": body,   "size_pt": S.BULL_BODY,
                   "bold": False, "color": BODY}])

    _add(tf, [
        {"text": "Platform- and DS-agnostic framework.",
         "size_pt": S.TAKEAWAY, "bold": True, "color": NAVY},
    ])


def _resize_col1_figures(slide):
    """Find the two Franka pictures in Col1 (top y around 34.5) and
    enlarge them vertically.  Keeps left/width unchanged."""
    fig_shapes = []
    for sh in slide.shapes:
        if sh.shape_type == 13:                                # PICTURE
            t_in = sh.top  / EMU
            l_in = sh.left / EMU
            if 33.0 < t_in < 36.0 and l_in < 11.0:
                fig_shapes.append(sh)
    fig_shapes.sort(key=lambda s: s.left)
    if not fig_shapes:
        print("  warn: no Col1 figures found to resize")
        return
    new_top = Emu(int(TARGET_FIG_TOP_IN * EMU))
    new_h   = Emu(int(TARGET_FIG_HEIGHT_IN * EMU))
    for sh in fig_shapes:
        old_t, old_h = sh.top / EMU, sh.height / EMU
        sh.top = new_top
        sh.height = new_h
        print(f"  resized {sh.name!r} (L={sh.left/EMU:.2f}): "
              f"T {old_t:.2f}->{TARGET_FIG_TOP_IN:.2f}, "
              f"H {old_h:.2f}->{TARGET_FIG_HEIGHT_IN:.2f}")


def _move_caption(cap, title):
    """Reposition Col1Caption to align with Col2Plan (T=38.92, H=4.78)
    and the column's left margin."""
    cap.left   = title.left
    cap.width  = title.width
    cap.top    = Emu(int(TARGET_CAP_TOP_IN * EMU))
    cap.height = Emu(int(TARGET_CAP_HEIGHT_IN * EMU))
    print(f"  moved Col1Caption to L={title.left/EMU:.2f} "
          f"T={TARGET_CAP_TOP_IN:.2f} H={TARGET_CAP_HEIGHT_IN:.2f}")


def _extend_separators(slide):
    """Stretch Col1Sep and Col2Sep so they reach the new column bottom."""
    for sh in slide.shapes:
        if sh.name in ("Col1Sep", "Col2Sep"):
            t_in = sh.top / EMU
            new_h = SEP_NEW_BOTTOM_IN - t_in
            old_h = sh.height / EMU
            sh.height = Emu(int(new_h * EMU))
            print(f"  extended {sh.name}: H {old_h:.2f}->{new_h:.2f}in")


def _align_other_plan_boxes(slide):
    """Align Col2Plan, Col3Plan, Col2Equation, Col3Box and the Col2/Col3
    figures with the same 4-row grid used by Col1.  All three columns end
    up bottom-flush at y=46.10 with identical row tops."""
    align_specs = [
        # (name, top_in, height_in)
        ("Col2Equation", TARGET_SETUP_TOP_IN, TARGET_SETUP_H_IN),
        ("Col3Box",      TARGET_SETUP_TOP_IN, TARGET_SETUP_H_IN),
        ("Col2Plan",     TARGET_CAP_TOP_IN,   TARGET_CAP_HEIGHT_IN),
        ("Col3Plan",     TARGET_CAP_TOP_IN,   TARGET_CAP_HEIGHT_IN),
    ]
    for nm, top_in, h_in in align_specs:
        sh = next((s for s in slide.shapes if s.name == nm), None)
        if sh is None:
            print(f"  warn: {nm} not found, skipping align")
            continue
        old_t, old_h = sh.top / EMU, sh.height / EMU
        sh.top    = Emu(int(top_in * EMU))
        sh.height = Emu(int(h_in   * EMU))
        # Disable auto-shrink so the shape height we set actually sticks
        if sh.has_text_frame:
            sh.text_frame.auto_size = MSO_AUTO_SIZE.NONE
            sh.text_frame.word_wrap = True
            # Center the equation rows vertically; leave the plan-row boxes
            # top-anchored so their bullet list flows from the top.
            if nm in ("Col2Equation", "Col3Box"):
                sh.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        print(f"  aligned {nm}: T {old_t:.2f}->{top_in:.2f}, "
              f"H {old_h:.2f}->{h_in:.2f}")

    # Col2Fig: explicitly resize to the matched bottom (B = 35.20 + 4.81 =
    # 40.01) so it bottom-aligns with Col3Fig and never overruns the plan row.
    sh = next((s for s in slide.shapes if s.name == "Col2Fig"), None)
    if sh is None:
        for s in slide.shapes:
            if s.shape_type == 13 and 11.0 < s.left / EMU < 22.0 \
               and 33.5 < s.top / EMU < 36.5:
                sh = s
                print(f"  fallback: matched Col2Fig by position ({s.name!r})")
                break
    if sh is not None:
        old_l, old_t = sh.left / EMU, sh.top / EMU
        old_w, old_h = sh.width / EMU, sh.height / EMU
        sh.left   = Emu(int(COL2_FIG_LEFT_IN   * EMU))
        sh.top    = Emu(int(TARGET_FIG_TOP_IN  * EMU))
        sh.width  = Emu(int(COL2_FIG_WIDTH_IN  * EMU))
        sh.height = Emu(int(COL2_FIG_HEIGHT_IN * EMU))
        print(f"  resized Col2Fig: L {old_l:.2f}->{COL2_FIG_LEFT_IN:.2f}, "
              f"T {old_t:.2f}->{TARGET_FIG_TOP_IN:.2f}, "
              f"W {old_w:.2f}->{COL2_FIG_WIDTH_IN:.2f}, "
              f"H {old_h:.2f}->{COL2_FIG_HEIGHT_IN:.2f}")

    # Col3Fig: also resize to full column width with the new (10x5) aspect,
    # so the workspace + time-series panels render large enough to be read.
    sh = next((s for s in slide.shapes if s.name == "Col3Fig"), None)
    if sh is None:
        for s in slide.shapes:
            if s.shape_type == 13 and s.left / EMU > 22.0 \
               and 33.5 < s.top / EMU < 36.5:
                sh = s
                print(f"  fallback: matched Col3Fig by position ({s.name!r})")
                break
    if sh is not None:
        old_l, old_t = sh.left / EMU, sh.top / EMU
        old_w, old_h = sh.width / EMU, sh.height / EMU
        sh.left   = Emu(int(COL3_FIG_LEFT_IN   * EMU))
        sh.top    = Emu(int(TARGET_FIG_TOP_IN  * EMU))
        sh.width  = Emu(int(COL3_FIG_WIDTH_IN  * EMU))
        sh.height = Emu(int(COL3_FIG_HEIGHT_IN * EMU))
        print(f"  resized Col3Fig: L {old_l:.2f}->{COL3_FIG_LEFT_IN:.2f}, "
              f"T {old_t:.2f}->{TARGET_FIG_TOP_IN:.2f}, "
              f"W {old_w:.2f}->{COL3_FIG_WIDTH_IN:.2f}, "
              f"H {old_h:.2f}->{COL3_FIG_HEIGHT_IN:.2f}")


def main():
    if not os.path.exists(BACKUP):
        shutil.copy2(POSTER, BACKUP)
        print(f"backup -> {BACKUP}")
    p = Presentation(POSTER)
    slide = p.slides[0]
    by_name = {sh.name: sh for sh in slide.shapes}
    for need in ("Col1Title", "Col1Caption"):
        if need not in by_name:
            raise RuntimeError(f"missing shape: {need}")
    title, cap = by_name["Col1Title"], by_name["Col1Caption"]
    for sh in (title, cap):
        sh.text_frame.auto_size = MSO_AUTO_SIZE.NONE
        sh.text_frame.word_wrap = True
    write_col1_title  (title.text_frame)
    setup = _ensure_col1_setup(slide, title)
    write_col1_setup  (setup.text_frame)
    _move_caption     (cap, title)
    write_col1_caption(cap.text_frame)
    _resize_col1_figures(slide)
    _align_other_plan_boxes(slide)
    _extend_separators  (slide)
    p.save(POSTER)
    print(f"saved -> {POSTER}")


if __name__ == "__main__":
    main()
