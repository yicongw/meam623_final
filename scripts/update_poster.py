"""Update Section 3 of meam623_poster.pptx with the LPV-DS results.

Changes
-------
- Col3Title: WIP -> Done badge (green), corrected formula, drop SEDS reference.
- Col3Box  : compact LPV-DS formula box (replaces 'SEDS -> LPV-DS').
- Inserts poster_fig3.png between Col3Box and Col3Plan.
- Col3Plan : results-and-take-aways bullet list (replaces 'To be shown').
"""

from __future__ import annotations
import copy
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

POSTER = Path(r"C:/Users/wangy/Desktop/meam623_poster.pptx")
FIG    = Path(r"C:/Users/wangy/Desktop/poster_fig3.png")
BACKUP = POSTER.with_suffix(".backup.pptx")

NAVY   = RGBColor(0x01, 0x1F, 0x5B)
GREEN  = RGBColor(0x2E, 0x7D, 0x32)
BODY   = RGBColor(0x33, 0x33, 0x33)
MUTED  = RGBColor(0x55, 0x55, 0x55)


def set_run(run, text, *, font="Arial", size=None, bold=None, color=None):
    run.text = text
    run.font.name = font
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def clear_text_frame(tf):
    """Remove all paragraphs except the first; clear that one of runs."""
    p_elems = tf._txBody.findall(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}p"
    )
    for p in p_elems[1:]:
        tf._txBody.remove(p)
    first = tf.paragraphs[0]
    for r in list(first.runs):
        first._p.remove(r._r)


def add_para(tf, *, alignment=None):
    p = tf.add_paragraph()
    if alignment is not None:
        p.alignment = alignment
    return p


def first_para(tf, *, alignment=None):
    p = tf.paragraphs[0]
    if alignment is not None:
        p.alignment = alignment
    return p


def main():
    if not POSTER.exists():
        raise FileNotFoundError(POSTER)
    if not FIG.exists():
        raise FileNotFoundError(FIG)

    if not BACKUP.exists():
        shutil.copy2(POSTER, BACKUP)
        print(f"[backup] -> {BACKUP}")

    prs = Presentation(str(POSTER))
    slide = prs.slides[0]

    shapes_by_name = {sh.name: sh for sh in slide.shapes}
    col3_title = shapes_by_name["Col3Title"]
    col3_box   = shapes_by_name["Col3Box"]
    col3_plan  = shapes_by_name["Col3Plan"]

    # -------------------------------------------------------------------- #
    # 1. Col3Title: 'Done' badge + corrected one-line summary.             #
    # -------------------------------------------------------------------- #
    tf = col3_title.text_frame
    clear_text_frame(tf)

    p0 = first_para(tf)
    r = p0.add_run()
    set_run(r, "3. Swapping DS Planners", size=40, bold=True, color=NAVY)
    r = p0.add_run()
    set_run(r, "  \u2713 Done", size=28, bold=False, color=GREEN)

    p1 = add_para(tf)
    r = p1.add_run()
    set_run(
        r,
        "Replace hand-designed linear attractor f(x) = \u2212A(x \u2212 x*) "
        "with a learning-based LPV-DS, keeping the QP / safety layer untouched.",
        size=32, bold=False, color=BODY,
    )

    # -------------------------------------------------------------------- #
    # 2. Col3Box: compact LPV-DS formula card.                             #
    # -------------------------------------------------------------------- #
    # Slim it down to make room for the figure.
    col3_box.top    = Inches(33.20)
    col3_box.height = Inches(1.30)

    tf = col3_box.text_frame
    clear_text_frame(tf)

    p0 = first_para(tf, alignment=PP_ALIGN.CENTER)
    r = p0.add_run()
    set_run(
        r,
        "f(x) = \u03A3\u2096 \u03B3\u2096(x) \u00B7 (A\u2096 x + b\u2096),"
        "   A\u2096 + A\u2096\u1D40 \u227A 0",
        size=32, bold=True, color=NAVY,
    )

    p1 = add_para(tf, alignment=PP_ALIGN.CENTER)
    r = p1.add_run()
    set_run(
        r,
        "Hand-drawn LPV-DS  \u00B7  K = 6 components  \u00B7  Hurwitz-projected",
        size=22, bold=False, color=MUTED,
    )

    # -------------------------------------------------------------------- #
    # 3. Insert poster_fig3.png between Col3Box and Col3Plan.              #
    # -------------------------------------------------------------------- #
    img_left   = col3_title.left           # = 22.20 in
    img_width  = col3_title.width          # = 9.62 in
    # 15:6 aspect ratio  ->  height = width * 6/15
    img_height = Emu(int(img_width * 6 / 15))
    img_top    = Inches(34.65)             # right under the slimmer Col3Box

    # Remove any earlier inserted figure (idempotent re-runs)
    for sh in list(slide.shapes):
        if sh.name == "Col3Fig":
            sp = sh._element
            sp.getparent().remove(sp)

    pic = slide.shapes.add_picture(
        str(FIG),
        left=img_left, top=img_top,
        width=img_width, height=img_height,
    )
    pic.name = "Col3Fig"

    # -------------------------------------------------------------------- #
    # 4. Col3Plan: actual results / take-aways.                            #
    # -------------------------------------------------------------------- #
    plan_top    = img_top + img_height + Inches(0.20)   # ~38.70 in
    # Stop above the bottom footer rectangle (which starts at top=43.64 in).
    plan_height = Inches(43.30) - plan_top
    col3_plan.top    = plan_top
    col3_plan.height = plan_height

    tf = col3_plan.text_frame
    clear_text_frame(tf)

    p0 = first_para(tf)
    r = p0.add_run()
    set_run(
        r,
        "Same QP & safety layer  \u2014  only f(x) is swapped.",
        size=28, bold=True, color=NAVY,
    )

    bullets = [
        ("\u2713  Safety preserved",
         "min d\u2092\u1D47\u209B = 7.9 cm,  min \u0393 = 15.9"),
        ("\u2713  Empirical passivity",
         "S(t) = T_kin + \u03BB V_pot stays bounded, "
         "\u224810\u00D7 lower transient than the linear DS"),
        ("\u2713  Hand-drawn LPV-DS converges",
         "final EE-to-target error 14.7 cm under viability constraints"),
    ]
    for head, sub in bullets:
        p = add_para(tf)
        r = p.add_run()
        set_run(r, head, size=28, bold=True, color=BODY)
        p2 = add_para(tf)
        r = p2.add_run()
        set_run(r, "    " + sub, size=22, bold=False, color=MUTED)

    p_take = add_para(tf)
    r = p_take.add_run()
    set_run(
        r,
        "Planner-agnostic VPP-TC: any GAS DS plugs in.",
        size=28, bold=True, color=NAVY,
    )

    # -------------------------------------------------------------------- #
    # Save                                                                  #
    # -------------------------------------------------------------------- #
    prs.save(str(POSTER))
    print(f"[save] -> {POSTER}")


if __name__ == "__main__":
    main()
