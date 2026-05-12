#!/usr/bin/env python3
"""Update Col2 (Model-Mismatch Robustness) — unified style.

Uses the shared style spec in scripts/poster_style.py so all 3 results
columns look like they were laid out by the same hand.
"""

from __future__ import annotations
import os, shutil, sys

from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE

_HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, _HERE)
import poster_style as S
from poster_style import NAVY, GREEN, DARK, BODY, MID, RED
from pptx.util import Pt

POSTER = r"C:\Users\wangy\Desktop\meam623_poster.pptx"
BACKUP = r"C:\Users\wangy\Desktop\meam623_poster_backup_before_col2.pptx"
FIG    = r"C:\meam623_finalproj\output\fig_poster_col2_timeseries.png"


def _clear(tf):
    p = tf.paragraphs[0]._p
    for sib in list(p.itersiblings()):
        sib.getparent().remove(sib)
    for r in list(p.findall(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}r")):
        p.remove(r)


def _add(tf, runs, *, align=None, first=False):
    para = tf.paragraphs[0] if first else tf.add_paragraph()
    if align is not None:
        para.alignment = align
    for spec in runs:
        run = para.add_run()
        run.text = spec["text"]
        f = run.font
        f.name = spec.get("name", S.FONT)
        f.size = Pt(spec["size_pt"])
        f.bold = spec.get("bold", False)
        f.italic = spec.get("italic", False)
        if spec.get("color") is not None:
            f.color.rgb = spec["color"]
        # LaTeX-style subscript / superscript via OOXML baseline attribute
        # (-25000 = subscript, +30000 = superscript)
        if spec.get("baseline") is not None:
            run._r.get_or_add_rPr().set("baseline", str(spec["baseline"]))


# Convenience builders for inline math fragments.
SUB = -25000   # baseline value for subscript
SUP =  30000   # baseline value for superscript


def _math(text, size_pt, *, color, bold=False):
    """Italic variable (e.g. d, x, t, S) with no decoration."""
    return {"text": text, "size_pt": size_pt, "italic": True,
            "bold": bold, "color": color}


def _sub(text, size_pt, *, color, bold=False):
    """Italic-subscript suffix (e.g. 'sc' in d_sc, 'fbk' in tau_fbk)."""
    return {"text": text, "size_pt": size_pt, "italic": True,
            "bold": bold, "color": color, "baseline": SUB}


def _sup(text, size_pt, *, color, bold=False):
    """Superscript suffix (e.g. * in x*)."""
    return {"text": text, "size_pt": size_pt, "italic": False,
            "bold": bold, "color": color, "baseline": SUP}


def _plain(text, size_pt, *, color, bold=False):
    return {"text": text, "size_pt": size_pt, "italic": False,
            "bold": bold, "color": color}


def write_col2_title(tf):
    _clear(tf)
    _add(tf, [
        {"text": "2. Model-Mismatch Robustness", "size_pt": S.TITLE,
         "bold": True, "color": NAVY},
        {"text": "  \u2713 Done", "size_pt": S.DONE,
         "bold": False, "color": GREEN},
    ], align=PP_ALIGN.LEFT, first=True)
    _add(tf, [
        {"text": "Each Panda link mass scaled by an i.i.d. factor "
                 "in [1\u2212p, 1+p], p \u2208 {20, 30, 40, 50}%. "
                 "We ablate the safety-fallback layer.",
         "size_pt": S.BODY_PT, "bold": False, "color": BODY},
    ])


def write_col2_equation(tf):
    """LaTeX-style equation: variables italic, operators / words upright."""
    _clear(tf)
    em = S.EQ_MAIN

    # Line 1:  Plant  M̃(q) q̈ = τ − h(q, q̇)   ≠   controller's M(q)
    # We keep things compact by writing h(q, q̇) as just "h" as before.
    _add(tf, [
        _plain("Plant  ", em, color=NAVY, bold=True),
        _math ("M\u0303", em, color=NAVY, bold=True),         # M̃
        _plain("(", em, color=NAVY, bold=True),
        _math ("q",  em, color=NAVY, bold=True),
        _plain(") ", em, color=NAVY, bold=True),
        _math ("q\u0308", em, color=NAVY, bold=True),         # q̈
        _plain(" = ", em, color=NAVY, bold=True),
        _math ("\u03c4", em, color=NAVY, bold=True),          # τ
        _plain(" \u2212 ", em, color=NAVY, bold=True),
        _math ("h", em, color=NAVY, bold=True),
        _plain("   \u2260   controller's  ", em, color=NAVY, bold=True),
        _math ("M", em, color=NAVY, bold=True),
        _plain("(", em, color=NAVY, bold=True),
        _math ("q", em, color=NAVY, bold=True),
        _plain(")", em, color=NAVY, bold=True),
    ], align=PP_ALIGN.CENTER, first=True)

    # Line 2:  Fallback  τ = ID(q, q̇, 0) − K · q̇   ·   96 cells × 6 seeds
    es = S.EQ_SUB
    _add(tf, [
        _plain("Fallback  ", es, color=MID),
        _math ("\u03c4", es, color=MID),                      # τ
        _plain(" = ID(", es, color=MID),
        _math ("q", es, color=MID),
        _plain(", ", es, color=MID),
        _math ("q\u0307", es, color=MID),                     # q̇
        _plain(", 0) \u2212 ", es, color=MID),
        _math ("K", es, color=MID),
        _plain(" \u00b7 ", es, color=MID),
        _math ("q\u0307", es, color=MID),
        _plain("    \u00b7    96 cells \u00d7 6 seeds", es, color=MID),
    ], align=PP_ALIGN.CENTER)


def write_col2_plan(tf):
    _clear(tf)
    bs = S.BULL_BODY  # short alias for math-fragment helpers

    _add(tf, [
        {"text": "With fallback: 0 / 48 self-collisions   "
                 "\u00b7   Without: 4 / 24",
         "size_pt": S.PLAN_HEAD, "bold": True, "color": NAVY},
    ], align=PP_ALIGN.LEFT, first=True)

    bullets = [
        # Bullet 1: WITH fallback -> safety holds
        ("\u2713  With fallback \u2014 safety holds.", DARK, [
            _plain("    0 / 48 self-collisions across ", bs, color=BODY),
            _math ("p", bs, color=BODY),
            _plain(" \u2208 {20\u201350}%; min ", bs, color=BODY),
            _math ("d", bs, color=BODY),
            _sub  ("sc", bs, color=BODY),
            _plain(" > 1 cm even at 50% mass mismatch.", bs, color=BODY),
        ]),
        # Bullet 2: WITHOUT fallback -> safety breaks
        ("\u2717  Without fallback \u2014 safety breaks.", RED, [
            _plain("    4 / 24 runs self-collide once OSQP becomes "
                   "infeasible \u2014 no safety net catches the arm.",
                   bs, color=BODY),
        ]),
        # Bullet 3: fallback brakes, doesn't drive to target
        ("\u2717  Fallback \u2260 convergence.", RED, [
            _plain("    ", bs, color=BODY),
            _math ("\u03c4", bs, color=BODY),
            _sub  ("fbk", bs, color=BODY),
            _plain(" only brakes (no target term); when it fires "
                   "\u226525% of steps the arm stalls 30\u201380 cm "
                   "from ", bs, color=BODY),
            _math ("x", bs, color=BODY),
            _sup  ("\u22c6", bs, color=BODY),
            _plain(". When OSQP stays feasible (13 / 24), the QP "
                   "drives the arm within 3 cm.", bs, color=BODY),
        ]),
    ]
    for header, hc, body_runs in bullets:
        _add(tf, [{"text": header, "size_pt": S.BULL_HEAD,
                   "bold": True, "color": hc}])
        _add(tf, body_runs)

    _add(tf, [
        {"text": "Fallback secures safety; QP feasibility delivers "
                 "convergence.",
         "size_pt": S.TAKEAWAY, "bold": True, "color": NAVY},
    ])


def replace_picture(slide, target_name, new_image_path,
                    *, expected_left_in=None, expected_top_in=None, tol_in=0.5):
    target = None
    for shape in slide.shapes:
        if shape.name == target_name:
            target = shape; break
    if target is None and expected_left_in is not None:
        EMU = 914400
        best, best_d = None, float("inf")
        for shape in slide.shapes:
            if shape.shape_type != 13:
                continue
            d = abs(shape.left/EMU - expected_left_in) + \
                abs(shape.top /EMU - expected_top_in)
            if d < best_d:
                best, best_d = shape, d
        if best is not None and best_d < tol_in * 2:
            print(f"  fallback: matched {target_name!r} by position "
                  f"(was {best.name!r}, dist={best_d:.2f}in)")
            target = best
    if target is None:
        raise RuntimeError(f"Picture {target_name!r} not found")
    left, top, w, h = target.left, target.top, target.width, target.height
    sp = target._element
    sp.getparent().remove(sp)
    pic = slide.shapes.add_picture(new_image_path, left, top, width=w, height=h)
    pic.name = target_name
    return pic


def main():
    if not os.path.exists(BACKUP):
        shutil.copy2(POSTER, BACKUP)
        print(f"backup -> {BACKUP}")
    p = Presentation(POSTER)
    slide = p.slides[0]
    by_name = {sh.name: sh for sh in slide.shapes}
    for need in ("Col2Title", "Col2Equation", "Col2Plan"):
        if need not in by_name:
            raise RuntimeError(f"missing shape: {need}")
    for nm in ("Col2Title", "Col2Equation", "Col2Plan"):
        tf = by_name[nm].text_frame
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.word_wrap = True
    write_col2_title    (by_name["Col2Title"   ].text_frame)
    write_col2_equation (by_name["Col2Equation"].text_frame)
    write_col2_plan     (by_name["Col2Plan"    ].text_frame)
    replace_picture     (slide, "Col2Fig", FIG,
                         expected_left_in=11.77, expected_top_in=34.68)
    p.save(POSTER)
    print(f"saved -> {POSTER}")


if __name__ == "__main__":
    main()
