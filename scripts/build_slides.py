#!/usr/bin/env python3
"""Build slide deck:  VPP-TC dynamics-mismatch ablation.

Output: output/vpptc_fallback_ablation.pptx  (16:9 widescreen)

Color palette: "Midnight Executive"
  primary   #1E2761  deep navy
  secondary #CADCFC  ice blue
  accent    #D7263D  coral red  (= "no fallback" in figures)
"""

from __future__ import annotations

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# ----------------------------- paths -----------------------------
_HERE        = os.path.dirname(os.path.abspath(__file__))
_PROJECT     = os.path.abspath(os.path.join(_HERE, os.pardir))
_OUT         = os.path.join(_PROJECT, "output")
PPTX_OUT     = os.path.join(_OUT, "vpptc_fallback_ablation.pptx")

FIG_COLL     = os.path.join(_OUT, "fig_hu_collisions.png")
FIG_QOVER    = os.path.join(_OUT, "fig_hu_qover.png")
FIG_SCDIST   = os.path.join(_OUT, "fig_hu_sc_dist.png")
FIG_METRICS  = os.path.join(_OUT, "fig_hu_metrics.png")

# ----------------------------- palette ---------------------------
NAVY   = RGBColor(0x1E, 0x27, 0x61)
NAVY_D = RGBColor(0x12, 0x18, 0x3D)   # darker shade
ICE    = RGBColor(0xCA, 0xDC, 0xFC)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_W = RGBColor(0xF5, 0xF7, 0xFB)
TEXT_D = RGBColor(0x1A, 0x1A, 0x2A)
TEXT_M = RGBColor(0x3A, 0x40, 0x60)   # mid-dark grey, readable on cream
ACCENT = RGBColor(0xD7, 0x26, 0x3D)

# ----------------------------- fonts -----------------------------
F_HEAD = "Calibri"
F_BODY = "Calibri"

# 16:9 widescreen
SLIDE_W_IN = 13.333
SLIDE_H_IN =  7.5


# =================================================================
def _set_bg(slide, prs, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    bg.shadow.inherit = False
    # send to back
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def _add_text(slide, x, y, w, h,
              text, *, font=F_BODY, size=18, bold=False, italic=False,
              color=TEXT_D, align=PP_ALIGN.LEFT,
              anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y),
                                   Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top  = tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = anchor

    paragraphs = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return tb


def _add_bullets(slide, x, y, w, h, items, *,
                 font=F_BODY, size=18, color=TEXT_D, bullet_color=NAVY,
                 line_spacing=1.30):
    """items: list of (header_str_or_None, body_str).
    If header is given, render it bold at +1pt; body slightly indented."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y),
                                   Inches(w), Inches(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Inches(0.04)
    tf.margin_top  = tf.margin_bottom = Inches(0.02)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP

    first = True
    for header, body in items:
        if header is not None:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = line_spacing
            r0 = p.add_run()
            r0.text = "■  "
            r0.font.name = font
            r0.font.size = Pt(size)
            r0.font.color.rgb = bullet_color
            r0.font.bold = True
            r1 = p.add_run()
            r1.text = header
            r1.font.name = font
            r1.font.size = Pt(size + 1)
            r1.font.bold = True
            r1.font.color.rgb = color
        if body:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = line_spacing
            r = p.add_run()
            r.text = ("    " if header else "") + body
            r.font.name = font
            r.font.size = Pt(size)
            r.font.color.rgb = color
    return tb


def _stat_card(slide, x, y, w, h,
               value, label, *, value_color=NAVY, label_color=TEXT_M,
               value_size=44, label_size=13, fill=NEAR_W,
               value_frac=0.50):
    """Rounded-rectangle stat card.

    `value_frac` controls the vertical share of the card given to the big
    value; the label gets the remainder (with a small gap).  Tune this
    when label has many lines so the bottom rows don't get clipped.
    """
    bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(x), Inches(y), Inches(w), Inches(h))
    bg.fill.solid()
    bg.fill.fore_color.rgb = fill
    bg.line.color.rgb = ICE
    bg.line.width = Pt(0.75)
    bg.shadow.inherit = False

    val_h   = h * value_frac
    label_y = y + val_h + 0.04
    label_h = h - val_h - 0.10

    _add_text(slide, x, y + 0.10, w, val_h - 0.05,
              value, font=F_HEAD, size=value_size, bold=True,
              color=value_color, align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.MIDDLE)
    _add_text(slide, x, label_y, w, label_h,
              label, font=F_BODY, size=label_size,
              color=label_color, align=PP_ALIGN.CENTER,
              anchor=MSO_ANCHOR.TOP, line_spacing=1.15)


def _slide_header(slide, prs, eyebrow, title, *,
                  on_dark=False):
    """Header band at top: eyebrow + title, no horizontal accent line."""
    eb_color = ICE if on_dark else NAVY
    title_color = WHITE if on_dark else NAVY_D
    _add_text(slide, 0.55, 0.34, 12.2, 0.36,
              eyebrow.upper(), font=F_HEAD, size=12, bold=True,
              color=eb_color, align=PP_ALIGN.LEFT)
    _add_text(slide, 0.55, 0.66, 12.2, 0.85,
              title, font=F_HEAD, size=30, bold=True,
              color=title_color, align=PP_ALIGN.LEFT,
              line_spacing=1.10)


def _footer(slide, prs, page_str, *, on_dark=False):
    color = ICE if on_dark else TEXT_M
    _add_text(slide, 0.55, 7.05, 9.0, 0.32,
              "VPP-TC dynamics-mismatch ablation  ·  Yicong Wang  ·  MEAM 623",
              font=F_BODY, size=9, color=color, align=PP_ALIGN.LEFT)
    _add_text(slide, 11.5, 7.05, 1.3, 0.32,
              page_str, font=F_BODY, size=9, color=color,
              align=PP_ALIGN.RIGHT)


# =================================================================
# Slide 1 — Title
# =================================================================
def slide_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_bg(s, prs, NAVY_D)

    # Decorative accent block (left side bar)
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             0, 0, Inches(0.42), prs.slide_height)
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background(); bar.shadow.inherit = False

    # Eyebrow
    _add_text(s, 1.05, 1.95, 11.0, 0.40,
              "MEAM 623  ·  FINAL-PROJECT EXPERIMENT",
              font=F_HEAD, size=14, bold=True, color=ICE)
    # Title
    _add_text(s, 1.05, 2.40, 11.5, 1.7,
              "Where does VPP-TC's robustness come from\nunder dynamics mismatch?",
              font=F_HEAD, size=40, bold=True, color=WHITE,
              line_spacing=1.10)
    # Sub
    _add_text(s, 1.05, 4.45, 11.5, 0.7,
              "An ablation of the gravity-comp + velocity-damping fallback layer.",
              font=F_BODY, size=20, italic=True, color=ICE)

    # Author + key numbers strip
    _add_text(s, 1.05, 6.35, 6.0, 0.40,
              "Yicong Wang", font=F_HEAD, size=18, bold=True, color=WHITE)
    _add_text(s, 1.05, 6.78, 6.0, 0.34,
              "Department of MEAM, University of Pennsylvania",
              font=F_BODY, size=13, color=ICE)

    _stat_card(s, 8.4, 5.85, 1.55, 1.45,
               "96", "simulated\ncells",
               value_color=ACCENT, label_color=ICE,
               value_size=32, label_size=11, fill=NAVY,
               value_frac=0.55)
    _stat_card(s, 10.05, 5.85, 1.55, 1.45,
               "0–50%", "link-mass\nperturbation",
               value_color=WHITE, label_color=ICE,
               value_size=22, label_size=11, fill=NAVY,
               value_frac=0.55)
    _stat_card(s, 11.7, 5.85, 1.55, 1.45,
               "2 modes", "fallback\nablation",
               value_color=WHITE, label_color=ICE,
               value_size=18, label_size=11, fill=NAVY,
               value_frac=0.55)


# =================================================================
# Slide 2 — Motivation / question
# =================================================================
def slide_question(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, prs, NEAR_W)
    _slide_header(s, prs, "Motivation",
                  "VPP-TC stacks several safety layers — which one matters?")

    # Left column: bullet list of safety layers
    _add_bullets(s, 0.55, 1.85, 6.5, 4.5, [
        ("Viability acceleration bound",
         "QP must keep q̈ inside a box that guarantees a future braking trajectory exists."),
        ("Self-collision Γ-constraint",
         "Learned predictor; QP enforces Γ̇ + ηΓ ≥ 0 to keep predicted distance positive."),
        ("Reactive SDF evasion",
         "Bypasses the QP when a real obstacle gets closer than a threshold."),
        ("Gravity-comp + damping  fallback",
         "When the QP itself is infeasible, send τ = g(q) − K·q̇ to bring the arm to rest passively."),
    ], size=15.5)

    # Right column: question card  (taller + a touch wider)
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(7.45), Inches(1.85),
                              Inches(5.45), Inches(4.85))
    card.fill.solid(); card.fill.fore_color.rgb = NAVY
    card.line.fill.background(); card.shadow.inherit = False

    _add_text(s, 7.65, 2.00, 5.05, 0.42,
              "THE QUESTION", font=F_HEAD, size=12, bold=True, color=ICE)
    _add_text(s, 7.65, 2.50, 5.10, 3.10,
              "When the simulator runs perturbed link masses but the "
              "controller still uses the nominal model, which layer of "
              "VPP-TC actually prevents self-collision?",
              font=F_HEAD, size=18, bold=True, color=WHITE,
              line_spacing=1.30)
    _add_text(s, 7.65, 5.75, 5.10, 0.85,
              "We compare two controller variants under the same "
              "PyBullet perturbations to isolate the fallback's role.",
              font=F_BODY, size=13, italic=True, color=ICE,
              line_spacing=1.25)
    _footer(s, prs, page)


# =================================================================
# Slide 3 — Experimental setup
# =================================================================
def slide_setup(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, prs, NEAR_W)
    _slide_header(s, prs, "Setup", "Proposal-aligned protocol  ·  96-cell sweep")

    # Left: protocol bullets
    _add_bullets(s, 0.55, 1.85, 6.5, 4.7, [
        ("Sim ≠ Controller",
         "PyBullet integrates with link masses scaled by 1 ± p (p ∈ {20, 30, 40, 50}%); "
         "controller's nominal-mass M(q) and τ_id are queried from a separate restored copy."),
        ("Apples-to-apples seeds",
         "Same RNG seed → same per-link perturbation pattern, so each (mode₁, mode₂) pair "
         "sees an identical robot."),
        ("Two ablation toggles",
         "(i) reactive-evasion ON / OFF;  (ii) safety fallback ON / OFF."),
        ("Speed",
         "duration = 3 s, 3 parallel workers — entire sweep finishes in ~19 min."),
    ], size=14.5)

    # Right: comparison table card
    card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(7.55), Inches(1.85),
                              Inches(5.30), Inches(4.70))
    card.fill.solid(); card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = ICE; card.line.width = Pt(1.0)
    card.shadow.inherit = False

    _add_text(s, 7.75, 2.00, 5.0, 0.45,
              "TWO MODES COMPARED", font=F_HEAD, size=12, bold=True, color=NAVY)

    # With fallback
    _add_text(s, 7.75, 2.55, 5.0, 0.40,
              "With fallback   (default VPP-TC)",
              font=F_HEAD, size=15, bold=True, color=NAVY)
    _add_text(s, 7.75, 2.95, 5.0, 1.0,
              "On QP infeasibility →\n   τ = ID(q, q̇, 0) − 5·q̇,   clipped to ±τ_max",
              font="Consolas", size=11, color=TEXT_D, line_spacing=1.30)

    # Divider
    div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(7.75), Inches(4.10),
                             Inches(4.85), Inches(0.012))
    div.fill.solid(); div.fill.fore_color.rgb = ICE
    div.line.fill.background(); div.shadow.inherit = False

    # No fallback
    _add_text(s, 7.75, 4.25, 5.0, 0.40,
              "No fallback   (ablation)",
              font=F_HEAD, size=15, bold=True, color=ACCENT)
    _add_text(s, 7.75, 4.65, 5.0, 1.0,
              "On QP infeasibility →\n   τ = clip(OSQP.u.value, ±τ_max);\n"
              "   τ = 0 if u.value is None",
              font="Consolas", size=11, color=TEXT_D, line_spacing=1.30)

    _footer(s, prs, page)


# =================================================================
# Slide 4 — Result: collision counts
# =================================================================
def slide_result_collisions(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, prs, NEAR_W)
    _slide_header(s, prs, "Result 1",
                  "The fallback eliminates self-collisions at every perturbation level")

    # Big stat callouts (top row)
    _stat_card(s, 0.55, 1.80, 2.85, 1.30,
               "0 / 24",
               "self-collisions, with fallback\n(reactive evasion ON)",
               value_color=NAVY, value_size=34, label_size=11)
    _stat_card(s, 3.55, 1.80, 2.85, 1.30,
               "4 / 24",
               "self-collisions, no fallback\n(reactive evasion ON)",
               value_color=ACCENT, value_size=34, label_size=11)
    _stat_card(s, 6.55, 1.80, 2.85, 1.30,
               "+1 cm",
               "mean min self-coll. distance\nadvantage (with vs without)",
               value_color=NAVY, value_size=30, label_size=11)
    _stat_card(s, 9.55, 1.80, 3.30, 1.30,
               "every pct",
               "no-fallback collides at\n20%, 30%, 40%, 50%",
               value_color=ACCENT, value_size=24, label_size=11)

    # Figure
    s.shapes.add_picture(FIG_COLL,
                          Inches(0.55), Inches(3.30),
                          width=Inches(8.4))

    # Right callout column
    _add_text(s, 9.30, 3.40, 3.6, 0.35,
              "WHAT TO SEE", font=F_HEAD, size=11, bold=True, color=NAVY)
    _add_bullets(s, 9.30, 3.65, 3.65, 3.0, [
        (None, "Blue bars (with fallback) are zero across all eight (pct, reactive) cells."),
        (None, "Red bars appear at every perturbation level when reactive evasion is on — and only one cell when it's off."),
        (None, "Reactive evasion + no fallback is the worst combination: evasion swings the arm into a bad config, then nothing brakes it."),
    ], size=11.5, line_spacing=1.30, bullet_color=ACCENT)

    _footer(s, prs, page)


# =================================================================
# Slide 5 — Result: q-overshoot reversal
# =================================================================
def slide_result_qover(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, prs, NEAR_W)
    _slide_header(s, prs, "Result 2",
                  "But the fallback pays for it in joint-position-limit overshoot")

    # Figure
    s.shapes.add_picture(FIG_QOVER,
                          Inches(0.55), Inches(1.85),
                          width=Inches(9.0))

    # Right column callouts
    _stat_card(s, 9.85, 1.85, 3.05, 1.20,
               "100%",
               "of 96 cells satisfy\nfallback q̈-over ≥ no-fallback",
               value_color=NAVY, value_size=30, label_size=10.5)

    _stat_card(s, 9.85, 3.20, 3.05, 1.20,
               "+2 to +5 mm",
               "mean q-limit overshoot\ndifference (per pct)",
               value_color=NAVY, value_size=18, label_size=10.5)

    _stat_card(s, 9.85, 4.55, 3.05, 1.20,
               "15.5 mm",
               "worst single-cell overshoot\n(fallback @ 50%)",
               value_color=ACCENT, value_size=22, label_size=10.5)

    _add_bullets(s, 9.85, 5.95, 3.10, 1.0, [
        (None, "Mechanism: the fallback's −K·q̇ term decelerates passively but does not actively respect joint-position bounds."),
    ], size=10.5, line_spacing=1.25, bullet_color=NAVY)

    _footer(s, prs, page)


# =================================================================
# Slide 6 — Mechanism
# =================================================================
def slide_mechanism(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, prs, NEAR_W)
    _slide_header(s, prs, "Mechanism",
                  "What the fallback layer is actually doing")

    # Three-step horizontal flow
    box_w, box_h = 4.0, 3.6
    y = 1.95
    xs = [0.55, 4.65, 8.75]
    titles = [
        "1.  QP becomes infeasible",
        "2.  With fallback",
        "3.  Without fallback",
    ]
    bodies = [
        "Under ±20–50% mass mismatch the inverse-dynamics term  τ_id = M_nom q̈ + …  diverges from the true value.\n\n"
        "The viability + Γ + torque constraints can no longer be jointly satisfied; OSQP returns "
        "infeasible or 'inaccurate'.",
        "Fallback fires:\n   τ  =  g(q)  −  K·q̇\n\nThis dissipates kinetic energy passively (τ·q̇ = −K‖q̇‖² ≤ 0); "
        "the arm decelerates and parks in a safe-ish configuration.\n\nq-limits may overshoot by mm, "
        "but the configuration stays out of self-collision.",
        "OSQP's last (inaccurate) solution is sent through directly. Joint accelerations remain near the "
        "viability bound — but the QP solution is wrong for the true dynamics.\n\n"
        "When reactive evasion has just swung the arm, momentum carries it into a self-collision configuration.",
    ]
    fills = [WHITE, NAVY, ACCENT]
    title_colors = [NAVY_D, WHITE, WHITE]
    body_colors = [TEXT_D, ICE, WHITE]

    for i, (x, title, body, fill, tc, bc) in enumerate(
            zip(xs, titles, bodies, fills, title_colors, body_colors)):
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                  Inches(x), Inches(y),
                                  Inches(box_w), Inches(box_h))
        card.fill.solid(); card.fill.fore_color.rgb = fill
        card.line.color.rgb = ICE
        card.line.width = Pt(0.75 if fill == WHITE else 0)
        card.shadow.inherit = False
        _add_text(s, x + 0.20, y + 0.20, box_w - 0.4, 0.55,
                  title, font=F_HEAD, size=17, bold=True, color=tc)
        _add_text(s, x + 0.20, y + 0.85, box_w - 0.4, box_h - 1.0,
                  body, font=F_BODY, size=12.5, color=bc, line_spacing=1.30)

    # Bottom takeaway band
    band = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(0.55), Inches(5.85),
                              Inches(12.20), Inches(0.95))
    band.fill.solid(); band.fill.fore_color.rgb = ICE
    band.line.fill.background(); band.shadow.inherit = False
    _add_text(s, 0.85, 5.97, 11.6, 0.75,
              "Take-away:  the fallback layer trades mm-level joint-tracking precision "
              "for a hard guarantee against self-collision when the QP cannot be trusted.",
              font=F_HEAD, size=15, italic=True, bold=True,
              color=NAVY_D, anchor=MSO_ANCHOR.MIDDLE)

    _footer(s, prs, page)


# =================================================================
# Slide 7 — Conclusion
# =================================================================
def slide_conclusion(prs, page):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(s, prs, NAVY_D)
    _slide_header(s, prs, "Conclusion",
                  "A layered defense is what makes VPP-TC robust",
                  on_dark=True)

    # Three stat cards — taller + slightly smaller value/label fonts so
    # all 3 lines of label fit comfortably.
    _stat_card(s, 0.55, 1.85, 4.0, 2.00,
               "0 / 48",
               "self-collisions in the\ndefault VPP-TC stack across\n0–50% mass perturbation",
               value_color=WHITE, label_color=ICE,
               value_size=42, label_size=12, fill=NAVY,
               value_frac=0.52)
    _stat_card(s, 4.70, 1.85, 4.0, 2.00,
               "4 / 24",
               "self-collisions appear when\nthe fallback layer alone is\nablated  (reactive ON, ε = 0)",
               value_color=ACCENT, label_color=ICE,
               value_size=42, label_size=12, fill=NAVY,
               value_frac=0.52)
    _stat_card(s, 8.85, 1.85, 4.0, 2.00,
               "+3 mm",
               "mean joint-limit overshoot\nthe fallback layer accepts\nas the price of safety",
               value_color=WHITE, label_color=ICE,
               value_size=42, label_size=12, fill=NAVY,
               value_frac=0.52)

    _add_text(s, 0.55, 4.15, 12.3, 0.40,
              "FINDINGS", font=F_HEAD, size=12, bold=True, color=ICE)
    _add_bullets(s, 0.55, 4.50, 12.3, 2.4, [
        ("State-based safety is intrinsically robust to model error.",
         "Γ depends only on (q, q̇) and the viability bound's structural conservatism absorbs ≤10% mass "
         "mismatch with no auxiliary mechanism — fallback never even fires."),
        ("The fallback layer carries the load at high uncertainty.",
         "From 20% upward the QP is regularly infeasible. Removing the fallback exposes the only "
         "self-collisions observed in the entire sweep."),
        ("ε-margin (proposal Section IV-c) does not help in this scenario.",
         "Mismatch shows up in τ_id, not at the viability bound — shrinking the q̈ box doesn't address "
         "the failure mode (verified separately, 60-cell sweep)."),
    ], size=14.0, color=WHITE, bullet_color=ACCENT, line_spacing=1.35)

    _footer(s, prs, page, on_dark=True)


# =================================================================
def main():
    prs = Presentation()
    prs.slide_width  = Inches(SLIDE_W_IN)
    prs.slide_height = Inches(SLIDE_H_IN)

    slide_title             (prs)
    slide_question          (prs, "2 / 7")
    slide_setup             (prs, "3 / 7")
    slide_result_collisions (prs, "4 / 7")
    slide_result_qover      (prs, "5 / 7")
    slide_mechanism         (prs, "6 / 7")
    slide_conclusion        (prs, "7 / 7")

    prs.save(PPTX_OUT)
    print(f"saved -> {PPTX_OUT}")


if __name__ == "__main__":
    main()
